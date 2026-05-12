#!/usr/bin/env python3
"""Bundle PNG slides + extracted speaker notes into a real .pptx (16:9 fullbleed)."""
import sys, html, re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

PPT_DIR = Path(__file__).resolve().parent.parent

HTML_PATH = sys.argv[1] if len(sys.argv) > 1 else str(PPT_DIR / 'index.html')
PNG_DIR   = sys.argv[2] if len(sys.argv) > 2 else str(PPT_DIR / 'ppt-png')
OUT_PATH  = sys.argv[3] if len(sys.argv) > 3 else str(PPT_DIR / 'jcli-thesis.pptx')

def clean(s):
    s = html.unescape(s)
    s = re.sub(r'\s+', ' ', s).strip()
    return s

def extract_notes(slide_el):
    """Pull all <p> from <aside class="notes">, return joined text."""
    notes = slide_el.select_one('aside.notes, .notes')
    if not notes:
        return ""
    paragraphs = []
    for p in notes.select('p'):
        t = clean(p.get_text())
        if t:
            paragraphs.append(t)
    return '\n\n'.join(paragraphs)

def main():
    html_content = Path(HTML_PATH).read_text(encoding='utf-8')
    soup = BeautifulSoup(html_content, 'lxml')
    slides_html = soup.select('.deck > section.slide')

    pngs = sorted(Path(PNG_DIR).glob('slide_*.png'))
    print(f"Found {len(slides_html)} slides in HTML, {len(pngs)} PNGs in {PNG_DIR}")
    if len(pngs) != len(slides_html):
        print(f"  ! mismatch — using min({len(slides_html)}, {len(pngs)})")
    n = min(len(pngs), len(slides_html))

    prs = Presentation()
    # 16:9 widescreen at 13.333 x 7.5 inches (matches PPT default widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for i in range(n):
        slide = prs.slides.add_slide(blank_layout)
        png_path = pngs[i]

        # Add picture filling the entire slide
        slide.shapes.add_picture(
            str(png_path),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Add speaker notes
        notes_text = extract_notes(slides_html[i])
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text

        if (i + 1) % 5 == 0 or i == n - 1:
            print(f"  + slide {i+1}/{n} ({png_path.name})")

    prs.save(OUT_PATH)
    size_kb = Path(OUT_PATH).stat().st_size / 1024
    print(f"\n✔ Saved {n} slides to {OUT_PATH} ({size_kb:.1f} KB)")
    print(f"  • Every slide is a fullbleed PNG of the HTML render")
    print(f"  • Speaker notes (逐字稿) embedded in 'View > Notes Page'")


if __name__ == '__main__':
    main()
